from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Formatierungskonstanten – analog zu slides.css

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

_BLUE_PRIMARY = RGBColor(0x25, 0x63, 0xEB)
_GREEN_PRIMARY = RGBColor(0x10, 0xB9, 0x81)
_GREEN_TITLE = RGBColor(0x06, 0x5F, 0x46)

_TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
_TEXT_SECONDARY = RGBColor(0x4B, 0x55, 0x63)
_TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)

_EXAMPLE_TEXT = RGBColor(0x78, 0x35, 0x0F)
_EXAMPLE_LABEL = RGBColor(0xB4, 0x53, 0x09)

_TITLE_SIZE_TITLE_SLIDE = Pt(44)
_TITLE_SIZE_CONTENT = Pt(32)
_TITLE_SIZE_CLOSING = Pt(36)
_BULLET_SIZE = Pt(18)
_BULLET_SIZE_TITLE = Pt(22)
_BULLET_SIZE_CLOSING = Pt(18)
_EXAMPLE_SIZE = Pt(16)
_FOOTER_SIZE = Pt(10)

_BULLET_SPACE_AFTER = Pt(6)
_EXAMPLE_SPACE_BEFORE = Pt(14)
_ACCENT_BAR_THICKNESS = Inches(0.1)

_MARGIN_LEFT = Inches(1.0)
_MARGIN_TOP = Inches(0.8)
_CONTENT_WIDTH = Inches(11.333)


# ---------------------------------------------------------------------------
# Hilfsfunktionen
# ---------------------------------------------------------------------------


# Zeichnet einen farbigen Akzentbalken oben oder links auf die Folie
def _add_accent_bar(slide: Any, color: RGBColor, *, top: bool = False, left: bool = False) -> None:
    if top:
        shape = slide.shapes.add_shape(1, Emu(0), Emu(0), _SLIDE_WIDTH, _ACCENT_BAR_THICKNESS)
    elif left:
        shape = slide.shapes.add_shape(1, Emu(0), Emu(0), _ACCENT_BAR_THICKNESS, _SLIDE_HEIGHT)
    else:
        return

    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# Erzeugt ein Textfeld an der angegebenen Position
def _add_textbox(slide: Any, left: Any, top: Any, width: Any, height: Any) -> Any:
    return slide.shapes.add_textbox(left, top, width, height)


# Fügt einen Paragraph mit optionalem Bullet-Marker (●) in den Textrahmen ein
def _add_bullet_paragraph(
    tf: Any,
    text: str,
    *,
    font_size: Any = _BULLET_SIZE,
    color: RGBColor = _TEXT_SECONDARY,
    bold: bool = False,
    show_marker: bool = True,
    alignment: int = PP_ALIGN.LEFT,
    space_after: Any = _BULLET_SPACE_AFTER,
) -> None:
    # Paragraph anlegen und Ausrichtung setzen
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after

    # Optionalen Bullet-Marker als eigenen Run voranstellen
    if show_marker:
        marker = p.add_run()
        marker.text = "●  "
        marker.font.size = Pt(8)
        marker.font.color.rgb = _BLUE_PRIMARY
        marker.font.bold = False

    # Eigentlichen Text formatieren
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold


# Fügt einen Example-Eintrag mit Präfix "Bsp.: " ein
def _add_example_paragraph(
    tf: Any,
    text: str,
    *,
    font_size: Any = _EXAMPLE_SIZE,
    alignment: int = PP_ALIGN.LEFT,
    first: bool = False,
) -> None:
    # Paragraph anlegen; beim ersten Example Extra-Abstand davor setzen
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = Pt(4)
    if first:
        p.space_before = _EXAMPLE_SPACE_BEFORE

    # Label "Bsp.: " fett in Akzentfarbe
    label_run = p.add_run()
    label_run.text = "Bsp.: "
    label_run.font.size = font_size
    label_run.font.color.rgb = _EXAMPLE_LABEL
    label_run.font.bold = True

    # Beispieltext kursiv
    text_run = p.add_run()
    text_run.text = text
    text_run.font.size = font_size
    text_run.font.color.rgb = _EXAMPLE_TEXT
    text_run.font.italic = True


# Fügt den Footer "AI-LV Assistant" unten rechts auf der Folie ein
def _add_footer(slide: Any, slide_number: int) -> None:
    # Textfeld unten rechts positionieren
    left = _SLIDE_WIDTH - Inches(3.0)
    top = _SLIDE_HEIGHT - Inches(0.55)
    txBox = _add_textbox(slide, left, top, Inches(2.8), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Brand-Text rechtsbündig einfügen
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    brand = p.add_run()
    brand.text = "AI-LV Assistant"
    brand.font.size = _FOOTER_SIZE
    brand.font.color.rgb = _TEXT_MUTED
    brand.font.bold = False



# Einziger öffentlicher Einstiegspunkt: erzeugt eine vollständige PPTX und gibt sie als Bytes zurück

def build_slides_pptx(slides: list[dict], topic: str) -> bytes:
    prs = _SlidePPTX(topic=topic)

    # Folien nach position sortieren, bevor sie verarbeitet werden
    for slide_data in sorted(slides, key=lambda s: s.get("position", 0)):
        prs.render_slide(slide_data)

    buffer = io.BytesIO()
    prs.prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


# Interne Klasse; kapselt die Präsentation und stellt Layout-Methoden bereit


class _SlidePPTX:

    # Initialisiert die Präsentation im 16:9-Format mit Blank-Layout
    def __init__(self, topic: str) -> None:
        self._topic = topic
        self._slide_count = 0
        self.prs = Presentation()
        self.prs.slide_width = _SLIDE_WIDTH
        self.prs.slide_height = _SLIDE_HEIGHT

    # Leitet anhand von slide_type an die passende Render-Methode weiter
    def render_slide(self, slide_data: dict[str, Any]) -> None:
        self._slide_count += 1
        slide_type = slide_data.get("slide_type", "content")
        if slide_type == "title":
            self._render_title(slide_data)
        elif slide_type == "closing":
            self._render_closing(slide_data)
        else:
            self._render_content(slide_data)

    # Titelfolie: blauer Akzent oben, großer zentrierter Titel, Bullets ohne Marker
    def _render_title(self, slide_data: dict[str, Any]) -> None:
        # Folie anlegen und blauen Akzentbalken oben einfügen
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        _add_accent_bar(slide, _BLUE_PRIMARY, top=True)

        # Titel zentriert und in Blau
        title_top = Inches(1.5)
        txBox = _add_textbox(slide, _MARGIN_LEFT, title_top, _CONTENT_WIDTH, Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data.get("title", "")
        run.font.size = _TITLE_SIZE_TITLE_SLIDE
        run.font.color.rgb = _BLUE_PRIMARY
        run.font.bold = True

        # Bullets zentriert ohne Marker unterhalb des Titels
        bullets: list[str] = slide_data.get("bullets") or []
        if bullets:
            body_top = title_top + Inches(2.2)
            body_box = _add_textbox(slide, _MARGIN_LEFT, body_top, _CONTENT_WIDTH, Inches(3.0))
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                bp = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
                bp.alignment = PP_ALIGN.CENTER
                bp.space_after = _BULLET_SPACE_AFTER
                run = bp.add_run()
                run.text = bullet
                run.font.size = _BULLET_SIZE_TITLE
                run.font.color.rgb = _TEXT_SECONDARY
                run.font.bold = False

        _add_footer(slide, self._slide_count)

    # Inhaltsfolie: blauer Akzent links, Titel mit Trennlinie, Bullets mit Marker
    def _render_content(self, slide_data: dict[str, Any]) -> None:
        # Folie anlegen und blauen Akzentbalken links einfügen
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        _add_accent_bar(slide, _BLUE_PRIMARY, left=True)

        # Titel linksbündig
        title_left = _MARGIN_LEFT + Inches(0.2)
        txBox = _add_textbox(slide, title_left, _MARGIN_TOP, _CONTENT_WIDTH, Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = slide_data.get("title", "")
        run.font.size = _TITLE_SIZE_CONTENT
        run.font.color.rgb = _TEXT_PRIMARY
        run.font.bold = True

        # Trennlinie unter dem Titel
        line_top = _MARGIN_TOP + Inches(0.85)
        line = slide.shapes.add_shape(1, title_left, line_top, Inches(6.0), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xDB, 0xE5, 0xF9)
        line.line.fill.background()

        # Bullets mit Marker und optionale Examples in den Body-Bereich schreiben
        bullets: list[str] = slide_data.get("bullets") or []
        examples: list[str] = slide_data.get("examples") or []
        if bullets or examples:
            body_top = line_top + Inches(0.35)
            body_box = _add_textbox(
                slide, title_left, body_top, _CONTENT_WIDTH, _SLIDE_HEIGHT - body_top - Inches(0.7),
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for bullet in bullets:
                _add_bullet_paragraph(body_tf, bullet, font_size=_BULLET_SIZE, color=_TEXT_SECONDARY, show_marker=True)
            for i, example in enumerate(examples):
                _add_example_paragraph(body_tf, example, first=(i == 0))

        _add_footer(slide, self._slide_count)

    # Abschlussfolie: grüner Akzent links, grüner Titel zentriert, Bullets ohne Marker
    def _render_closing(self, slide_data: dict[str, Any]) -> None:
        # Folie anlegen und grünen Akzentbalken links einfügen
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        _add_accent_bar(slide, _GREEN_PRIMARY, left=True)

        # Titel zentriert in Closing-Grün
        title_top = Inches(1.8)
        txBox = _add_textbox(slide, _MARGIN_LEFT, title_top, _CONTENT_WIDTH, Inches(1.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data.get("title", "")
        run.font.size = _TITLE_SIZE_CLOSING
        run.font.color.rgb = _GREEN_TITLE
        run.font.bold = True

        # Bullets zentriert ohne Marker und optionale Examples
        bullets: list[str] = slide_data.get("bullets") or []
        examples: list[str] = slide_data.get("examples") or []
        if bullets or examples:
            body_top = title_top + Inches(1.8)
            body_box = _add_textbox(slide, _MARGIN_LEFT, body_top, _CONTENT_WIDTH, Inches(3.0))
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for bullet in bullets:
                _add_bullet_paragraph(
                    body_tf, bullet,
                    font_size=_BULLET_SIZE_CLOSING,
                    color=_TEXT_SECONDARY,
                    show_marker=False,
                    alignment=PP_ALIGN.CENTER,
                )
            for i, example in enumerate(examples):
                _add_example_paragraph(body_tf, example, alignment=PP_ALIGN.CENTER, first=(i == 0))

        _add_footer(slide, self._slide_count)
